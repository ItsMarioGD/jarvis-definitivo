#!/usr/bin/env python3
"""
jarvis_generator.py - Modulo de generacion universal de JARVIS v2
Genera: documentos, imagenes, diagramas, modelos 3D, planos, codigos, etc.
"""
import os, re, math, json, time, random
from datetime import datetime

OUTPUT_BASE = os.path.join(os.path.expanduser("~"), "Descargas", "JARVIS", "Generaciones")

def _out(subfolder, filename):
    path = os.path.join(OUTPUT_BASE, subfolder, datetime.now().strftime("%Y-%m-%d"))
    os.makedirs(path, exist_ok=True)
    return os.path.join(path, filename)

def _safe_name(text, maxlen=40):
    return re.sub(r"[^\w\s-]", "", text).strip().replace(" ", "_")[:maxlen]


class JarvisGenerator:
    def __init__(self, log=print):
        self.log = log
        self._ollama_url = "http://localhost:11434/api/generate"
        self._llm_model = "llama3.2:1b"

    def _llm_generate(self, system_prompt, user_prompt, max_tokens=1024):
        """Llama 3 via Ollama para contenido inteligente"""
        try:
            import requests as _req
            resp = _req.post(self._ollama_url, json={
                "model": self._llm_model,
                "system": system_prompt,
                "prompt": user_prompt,
                "stream": False,
                "options": {"num_predict": max_tokens, "temperature": 0.7}
            }, timeout=60)
            resp.raise_for_status()
            return resp.json().get("response", "").strip()
        except Exception as e:
            self.log(f"Error Llama 3: {e}")
            return ""

    def _llm_sections(self, prompt, n_sections=5):
        """Genera secciones de contenido usando Llama 3"""
        sys = (
            "Eres un escritor profesional experto. Genera contenido ORIGINAL, detallado y de alta calidad. "
            "Responde SOLO con un JSON valido: [{\"title\": \"...\", \"body\": \"...\"}, ...] "
            "Sin markdown, sin explicaciones, solo el JSON puro."
        )
        user = (
            f"Genera exactamente {n_sections} secciones de contenido para un documento sobre: \"{prompt[:200]}\"\n"
            "Cada seccion debe tener un titulo corto y un cuerpo de 2-4 parrafos bien escritos.\n"
            "Formato: [{\"title\": \"Titulo\", \"body\": \"Contenido aqui...\"}]"
        )
        raw = self._llm_generate(sys, user, 2048)
        if raw:
            import json as _json, re as _re
            try:
                raw_clean = raw.strip()
                # Quitar bloques de codigo markdown
                code_match = _re.search(r'```(?:json)?\s*\n?(.*?)\n?```', raw_clean, _re.DOTALL)
                if code_match:
                    raw_clean = code_match.group(1).strip()
                # Buscar array JSON
                bracket_start = raw_clean.find('[')
                bracket_end = raw_clean.rfind(']')
                if bracket_start >= 0 and bracket_end > bracket_start:
                    raw_clean = raw_clean[bracket_start:bracket_end + 1]
                sections = _json.loads(raw_clean)
                if isinstance(sections, list) and len(sections) > 0:
                    # Filtrar solo title+body, ignorar campos extra
                    clean = []
                    for s in sections:
                        if isinstance(s, dict):
                            title = s.get("title", "")
                            body = s.get("body", "")
                            if title and body:
                                clean.append({"title": title, "body": body})
                    if clean:
                        return clean
            except Exception as e:
                self.log(f"Error parseando JSON Llama 3: {e}")
        return None

    def generate(self, request_text):
        t = request_text.lower().strip()

        # Detectar tipo por palabras clave (orden importa: especificos primero)
        if any(k in t for k in ["imagen", "foto", "dibujo", "ilustracion", "pic", "drawing", "render", "paint"]):
            return self.gen_image(request_text)
        if any(k in t for k in ["diagrama", "flowchart", "flujo", "organigrama", "chart", "grafico"]):
            return self.gen_diagram(request_text)
        if any(k in t for k in ["modelo 3d", "model 3d", "stl", "objeto 3d", "3d model", "cubo", "esfera 3d", "piramide", "cilindro"]):
            return self.gen_3d_model(request_text)
        if any(k in t for k in ["plano", "blueprint", "layout", "plano arquitectonico", "casa", "oficina", "habitacion"]):
            return self.gen_blueprint(request_text)
        if any(k in t for k in ["powerpoint", "presentacion", "pptx", "slides", "diapositiva"]):
            return self.gen_pptx(request_text)
        if any(k in t for k in ["word", "docx", "contrato", "texto formateado", "carta", "oficio"]) and "documento" not in t:
            return self.gen_word(request_text)
        if any(k in t for k in ["excel", "hoja de calculo", "spreadsheet", "tabla de datos"]):
            return self.gen_excel(request_text)
        if any(k in t for k in ["codigo", "code", "script", "programa", "funcion", "clase", "python", "javascript", "html"]):
            return self.gen_code(request_text)
        if any(k in t for k in ["plan", "proyecto", "planificacion", "roadmap", "timeline", "cronograma", "fases"]):
            return self.gen_plan(request_text)
        if any(k in t for k in ["musica", "melodia", "cancion", "music", "song", "audio", "tono", "ringtone"]):
            return self.gen_music(request_text)
        if any(k in t for k in ["video", "animacion", "animation", "clip"]):
            return {"type": "video", "message": "Generacion de video requiere modelo externo. Use imagen o animacion como alternativa.", "prompt": request_text}

        # Default: documento
        return self.gen_document(request_text)

    # ═══════════════════════════════════════════════════════════════════════════
    # IMAGEN - Generacion real con Pollinations AI (Flux)
    # ═══════════════════════════════════════════════════════════════════════════
    def gen_image(self, prompt):
        try:
            import requests as _req
            from urllib.parse import quote

            enhanced = self._enhance_image_prompt(prompt)
            seed = sum(ord(c) for c in prompt) % 100000

            self.log(f"Generando imagen real con IA: {enhanced[:80]}...")

            url = f"https://image.pollinations.ai/p/{quote(enhanced)}?width=1024&height=1024&seed={seed}&model=flux&nologo=true"

            resp = None
            for attempt in range(3):
                try:
                    resp = _req.get(url, timeout=120)
                    ct = resp.headers.get("Content-Type", "")
                    if resp.status_code == 200 and "image" in ct:
                        break
                    self.log(f"Intento {attempt+1}: status={resp.status_code} ct={ct}")
                    resp = None
                    import time as _t; _t.sleep(2)
                except Exception as e:
                    self.log(f"Intento {attempt+1} error: {e}")
                    resp = None
                    import time as _t; _t.sleep(2)

            if resp is None or resp.status_code != 200 or "image" not in resp.headers.get("Content-Type", ""):
                return self._gen_image_procedural(prompt)

            fname = _safe_name(prompt, 30) + "_ai.png"
            path = _out("Imagenes", fname)
            with open(path, "wb") as f:
                f.write(resp.content)

            size_kb = len(resp.content) / 1024
            self.log(f"Imagen generada ({size_kb:.0f} KB): {path}")
            return {"type": "image", "path": path, "prompt": prompt, "size_kb": round(size_kb)}

        except Exception as e:
            self.log(f"Error generando imagen con Pollinations: {e}")
            return self._gen_image_procedural(prompt)

    def _enhance_image_prompt(self, prompt):
        low = prompt.lower().strip()
        if any(k in low for k in ["foto", "photo", "retrato", "portrait", "selfie"]):
            return f"high quality photograph, {prompt}, realistic lighting, sharp focus, professional photography, 8k"
        if any(k in low for k in ["paisaje", "landscape", "naturaleza", "mountain", "beach", "forest"]):
            return f"stunning landscape photograph, {prompt}, golden hour lighting, vivid colors, cinematic, 8k uhd"
        if any(k in low for k in ["ciudad", "city", "cyberpunk", "futur", "skyline", "arquitectura"]):
            return f"epic cityscape, {prompt}, dramatic lighting, neon glow, cinematic, ultra detailed, 8k"
        if any(k in low for k in ["gato", "cat", "perro", "dog", "animal", "mascota", "pet"]):
            return f"adorable animal photo, {prompt}, soft bokeh background, warm lighting, ultra detailed fur, 8k"
        if any(k in low for k in ["comida", "food", "plato", "dish", "receta", "recipe"]):
            return f"professional food photography, {prompt}, studio lighting, appetizing, sharp focus, 8k"
        if any(k in low for k in ["dibujo", "drawing", "sketch", "pencil", "lapiz", "carbon"]):
            return f"detailed pencil sketch, {prompt}, graphite on paper, fine lines, artistic, high contrast"
        if any(k in low for k in ["anime", "manga", "cartoon", "animacion"]):
            return f"beautiful anime illustration, {prompt}, vibrant colors, clean linework, studio ghibli style, 4k"
        if any(k in low for k in ["oil", "acrylic", "painting", "pintura", "oleo", "acrilico"]):
            return f"masterful oil painting, {prompt}, rich textures, dramatic chiaroscuro, museum quality"
        if any(k in low for k in ["3d", "render", "blender", "octane"]):
            return f"photorealistic 3d render, {prompt}, octane render, volumetric lighting, subsurface scattering, 8k"
        return f"{prompt}, high quality, detailed, professional, sharp focus, vivid colors"

    def _gen_image_procedural(self, prompt):
        """Fallback procedural si Pollinations falla"""
        try:
            from PIL import Image, ImageDraw, ImageFont
            w, h = 1024, 1024
            rng = random.Random(sum(ord(c) for c in prompt) % 100000)
            accent = (0, 180, 240)
            img = Image.new('RGB', (w, h), (10, 12, 20))
            draw = ImageDraw.Draw(img)
            for y in range(h):
                ratio = y / h
                r = int(10 * (1-ratio) + 25 * ratio)
                g = int(12 * (1-ratio) + 35 * ratio)
                b = int(20 * (1-ratio) + 55 * ratio)
                draw.line([(0, y), (w, y)], fill=(r, g, b))
            for _ in range(rng.randint(20, 50)):
                x1, y1 = rng.randint(-50, w), rng.randint(-50, h)
                x2, y2 = x1 + rng.randint(100, 400), y1 + rng.randint(100, 400)
                alpha = rng.randint(15, 45)
                cr = max(0, min(255, accent[0] + rng.randint(-40, 40)))
                cg = max(0, min(255, accent[1] + rng.randint(-40, 40)))
                cb = max(0, min(255, accent[2] + rng.randint(-40, 40)))
                if rng.random() > 0.5:
                    draw.ellipse([x1, y1, x2, y2], fill=(cr, cg, cb, alpha))
                else:
                    draw.rectangle([x1, y1, x2, y2], fill=(cr, cg, cb, alpha))
            try:
                font = ImageFont.truetype("arial.ttf", 28)
            except:
                font = ImageFont.load_default()
            draw.text((30, 30), "JARVIS AI", fill=accent, font=font)
            fname = _safe_name(prompt, 30) + "_fallback.png"
            path = _out("Imagenes", fname)
            img.save(path, "PNG")
            self.log(f"Imagen procedural (fallback): {path}")
            return {"type": "image", "path": path, "prompt": prompt, "fallback": True}
        except Exception as e:
            return {"type": "error", "message": str(e)}

    # ═══════════════════════════════════════════════════════════════════════════
    # DIAGRAMA
    # ═══════════════════════════════════════════════════════════════════════════
    def gen_diagram(self, prompt):
        try:
            import matplotlib
            matplotlib.use('Agg')
            import matplotlib.pyplot as plt
            import matplotlib.patches as mpatches
            from matplotlib.patches import FancyBboxPatch

            fig, ax = plt.subplots(1, 1, figsize=(14, 9), facecolor='#0a0e18')
            ax.set_facecolor('#0a0e18')
            ax.set_xlim(0, 100)
            ax.set_ylim(0, 100)
            ax.axis('off')

            t = prompt.lower()
            if any(k in t for k in ["flujo", "flowchart", "proceso"]):
                steps = ["Inicio", "Analisis", "Diseno", "Desarrollo", "Pruebas", "Despliegue", "Fin"]
                for i, step in enumerate(steps):
                    x = 12 + (i % 4) * 22
                    y = 78 - (i // 4) * 35
                    color = '#00c8ff' if i in [0, len(steps)-1] else '#ff6a00'
                    box = FancyBboxPatch((x-9, y-6), 18, 12, boxstyle="round,pad=0.5",
                                           facecolor='#0d1520', edgecolor=color, linewidth=2.5)
                    ax.add_patch(box)
                    ax.text(x, y, step, ha='center', va='center', color='white', fontsize=11, fontweight='bold')
                    if i < len(steps) - 1:
                        nx = 12 + ((i+1) % 4) * 22
                        ny = 78 - ((i+1) // 4) * 35
                        ax.annotate('', xy=(nx-9, ny), xytext=(x+9, y),
                                   arrowprops=dict(arrowstyle='->', color='#00c8ff', lw=2))
            elif any(k in t for k in ["organigrama", "organizacion", "empresa"]):
                boxes = [
                    ("CEO", 50, 85, '#00c8ff'),
                    ("CTO", 25, 65, '#ff6a00'), ("CFO", 75, 65, '#ff6a00'),
                    ("Dev", 15, 45, '#22cc66'), ("QA", 35, 45, '#22cc66'),
                    ("Finanzas", 65, 45, '#22cc66'), ("RRHH", 85, 45, '#22cc66'),
                ]
                for name, x, y, color in boxes:
                    box = FancyBboxPatch((x-8, y-5), 16, 10, boxstyle="round,pad=0.3",
                                           facecolor='#0d1520', edgecolor=color, linewidth=2)
                    ax.add_patch(box)
                    ax.text(x, y, name, ha='center', va='center', color='white', fontsize=10, fontweight='bold')
                # Conexiones
                ax.annotate('', xy=(25, 70), xytext=(50, 80), arrowprops=dict(arrowstyle='->', color='#4a6878', lw=1.5))
                ax.annotate('', xy=(75, 70), xytext=(50, 80), arrowprops=dict(arrowstyle='->', color='#4a6878', lw=1.5))
                for x in [15, 35]:
                    ax.annotate('', xy=(x, 50), xytext=(25, 60), arrowprops=dict(arrowstyle='->', color='#4a6878', lw=1))
                for x in [65, 85]:
                    ax.annotate('', xy=(x, 50), xytext=(75, 60), arrowprops=dict(arrowstyle='->', color='#4a6878', lw=1))
            else:
                nodes = ["Concepto", "Plan", "Ejecucion", "Resultado", "Revision"]
                angles = [math.radians(i * 360 / len(nodes) - 90) for i in range(len(nodes))]
                cx, cy, r = 50, 50, 30
                for i, (node, angle) in enumerate(zip(nodes, angles)):
                    x = cx + r * math.cos(angle)
                    y = cy + r * math.sin(angle)
                    circle = plt.Circle((x, y), 9, facecolor='#0d1520', edgecolor='#00c8ff', linewidth=2.5)
                    ax.add_patch(circle)
                    ax.text(x, y, node, ha='center', va='center', color='white', fontsize=10, fontweight='bold')
                    next_angle = angles[(i+1) % len(nodes)]
                    nx = cx + r * math.cos(next_angle)
                    ny = cy + r * math.sin(next_angle)
                    ax.annotate('', xy=(nx, ny), xytext=(x, y),
                               arrowprops=dict(arrowstyle='->', color='#ff6a00', lw=2))

            ax.set_title(prompt[:80], color='#00c8ff', fontsize=15, fontweight='bold', pad=20)

            fname = _safe_name(prompt, 30) + "_diagram.png"
            path = _out("Diagramas", fname)
            fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#0a0e18')
            plt.close(fig)
            self.log(f"Diagrama generado: {path}")
            return {"type": "diagram", "path": path, "prompt": prompt}

        except Exception as e:
            self.log(f"Error generando diagrama: {e}")
            return {"type": "error", "message": str(e)}

    # ═══════════════════════════════════════════════════════════════════════════
    # MODELO 3D
    # ═══════════════════════════════════════════════════════════════════════════
    def gen_3d_model(self, prompt):
        try:
            import numpy as np
            from stl import mesh

            t = prompt.lower()
            seed = sum(ord(c) for c in prompt) % 10000
            rng = random.Random(seed)

            if any(k in t for k in ["cubo", "cube", "caja", "box"]):
                v = np.array([
                    [-1,-1,-1],[1,-1,-1],[1,1,-1],[-1,1,-1],
                    [-1,-1,1],[1,-1,1],[1,1,1],[-1,1,1]
                ], dtype=np.float64) * 50
                f = np.array([
                    [0,3,1],[1,3,2],[4,5,7],[4,7,6],
                    [0,4,3],[3,4,7],[1,2,5],[2,6,5],
                    [0,1,4],[1,5,4],[2,3,6],[3,7,6]
                ])
            elif any(k in t for k in ["esfera", "sphere", "bola", "orb"]):
                lat, lon = 24, 24
                v = []
                for i in range(lat+1):
                    theta = math.pi * i / lat
                    for j in range(lon):
                        phi = 2 * math.pi * j / lon
                        x = 50 * math.sin(theta) * math.cos(phi)
                        y = 50 * math.sin(theta) * math.sin(phi)
                        z = 50 * math.cos(theta)
                        v.append([x, y, z])
                v = np.array(v)
                f = []
                for i in range(lat):
                    for j in range(lon):
                        p1 = i * lon + j
                        p2 = i * lon + (j+1) % lon
                        p3 = (i+1) * lon + (j+1) % lon
                        p4 = (i+1) * lon + j
                        f.append([p1, p2, p3])
                        f.append([p1, p3, p4])
                f = np.array(f)
            elif any(k in t for k in ["piramide", "pyramid", "conico"]):
                v = np.array([
                    [0, 60, 0],
                    [-40, -20, -40], [40, -20, -40],
                    [40, -20, 40], [-40, -20, 40]
                ], dtype=np.float64)
                f = np.array([[0,1,2],[0,2,3],[0,3,4],[0,4,1],[1,4,3],[1,3,2]])
            elif any(k in t for k in ["cilindro", "cylinder", "tubo"]):
                segs = 20
                v = []
                for i in range(segs):
                    angle = 2 * math.pi * i / segs
                    v.append([30*math.cos(angle), -40, 30*math.sin(angle)])
                    v.append([30*math.cos(angle), 40, 30*math.sin(angle)])
                v.append([0, -40, 0])  # bottom center
                v.append([0, 40, 0])   # top center
                v = np.array(v)
                f = []
                bc = len(v)-2; tc = len(v)-1
                for i in range(segs):
                    n = (i+1) % segs
                    f.append([i*2, n*2, bc])
                    f.append([i*2+1, n*2+1, tc])
                    f.append([i*2, i*2+1, n*2+1])
                    f.append([i*2, n*2+1, n*2])
                f = np.array(f)
            else:
                # Forma aleatoria
                nv = rng.randint(10, 35)
                v = np.array([[rng.uniform(-50,50) for _ in range(3)] for _ in range(nv)])
                faces_list = []
                for i in range(nv):
                    for j in range(i+1, min(i+4, nv)):
                        for k in range(j+1, min(j+3, nv)):
                            if rng.random() > 0.5:
                                faces_list.append([i, j, k])
                f = np.array(faces_list) if faces_list else np.array([[0,1,2]])

            num_faces = len(f)
            model = mesh.Mesh(np.zeros(num_faces, dtype=mesh.Mesh.dtype))
            for i, face in enumerate(f):
                for j in range(3):
                    model.vectors[i][j] = v[face[j]]

            fname = _safe_name(prompt, 30) + "_3d.stl"
            path = _out("Modelos3D", fname)
            model.save(path)
            self.log(f"Modelo 3D generado: {path} ({num_faces} caras)")
            return {"type": "3d_model", "path": path, "prompt": prompt, "faces": num_faces}

        except Exception as e:
            self.log(f"Error generando modelo 3D: {e}")
            return {"type": "error", "message": str(e)}

    # ═══════════════════════════════════════════════════════════════════════════
    # PLANO
    # ═══════════════════════════════════════════════════════════════════════════
    def gen_blueprint(self, prompt):
        try:
            import matplotlib
            matplotlib.use('Agg')
            import matplotlib.pyplot as plt
            import matplotlib.patches as patches

            fig, ax = plt.subplots(1, 1, figsize=(14, 10), facecolor='#0a1020')
            ax.set_facecolor('#0a1020')
            ax.set_xlim(0, 140)
            ax.set_ylim(0, 100)
            ax.set_aspect('equal')

            t = prompt.lower()
            if any(k in t for k in ["casa", "house", "vivienda", "apartamento", "piso"]):
                rooms = [
                    ("Sala", 10, 55, 40, 30), ("Cocina", 50, 55, 30, 30),
                    ("Dormitorio\nPrincipal", 80, 55, 40, 30), ("Bano", 120, 55, 15, 30),
                    ("Garaje", 10, 10, 35, 40), ("Patio", 50, 10, 50, 40),
                    ("Dormitorio 2", 100, 10, 25, 40), ("Estudio", 125, 10, 10, 40),
                ]
                edge = '#00c8ff'
            elif any(k in t for k in ["oficina", "office", "empresa", "cowork"]):
                rooms = [
                    ("Recepcion", 10, 65, 30, 25), ("Sala Reuniones", 40, 65, 35, 25),
                    ("Open Space", 75, 65, 55, 25), ("Director", 10, 25, 25, 35),
                    ("Servidores", 35, 25, 20, 35), ("Banos", 55, 25, 15, 35),
                    ("Cocina", 70, 25, 25, 35), ("Archivos", 95, 25, 35, 35),
                ]
                edge = '#ff6a00'
            else:
                rooms = []
                for i in range(4):
                    for j in range(3):
                        rooms.append((f"H{i*3+j+1}", 10+i*32, 10+j*28, 28, 24))
                edge = '#00c8ff'

            for name, x, y, w, h in rooms:
                rect = patches.Rectangle((x, y), w, h, linewidth=2.5, edgecolor=edge, facecolor='#0d1520')
                ax.add_patch(rect)
                ax.text(x + w/2, y + h/2, name, ha='center', va='center', color=edge, fontsize=9, fontweight='bold')
                ax.text(x + w/2, y + h/2 - 3.5, f"{w}x{h}", ha='center', va='center', color='#3a5a6a', fontsize=7)

            ax.set_title(prompt[:80], color='#00c8ff', fontsize=14, fontweight='bold', pad=15)
            ax.grid(True, color='#1a2535', linewidth=0.5)

            for x in range(0, 141, 10):
                ax.text(x, 2, str(x), ha='center', va='center', color='#2a4050', fontsize=7)
            for y in range(0, 101, 10):
                ax.text(2, y, str(y), ha='center', va='center', color='#2a4050', fontsize=7)

            fname = _safe_name(prompt, 30) + "_plano.png"
            path = _out("Planos", fname)
            fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#0a1020')
            plt.close(fig)
            self.log(f"Plano generado: {path}")
            return {"type": "blueprint", "path": path, "prompt": prompt}

        except Exception as e:
            self.log(f"Error generando plano: {e}")
            return {"type": "error", "message": str(e)}

    # ═══════════════════════════════════════════════════════════════════════════
    # DOCUMENTOS
    # ═══════════════════════════════════════════════════════════════════════════
    def gen_document(self, prompt):
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.colors import HexColor

            fname = _safe_name(prompt, 30) + ".pdf"
            path = _out("Documentos", fname)

            doc = SimpleDocTemplate(path, pagesize=A4, topMargin=1*inch, bottomMargin=1*inch)
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle('T', parent=styles['Title'], textColor=HexColor('#00c8ff'), fontSize=18, spaceAfter=20)
            body_style = ParagraphStyle('B', parent=styles['Normal'], textColor=HexColor('#333333'), fontSize=11, spaceAfter=10, leading=16)
            head_style = ParagraphStyle('H', parent=styles['Heading2'], textColor=HexColor('#ff6a00'), fontSize=14, spaceAfter=6)

            elements = []
            elements.append(Paragraph(prompt[:100], title_style))
            elements.append(Spacer(1, 12))
            elements.append(Paragraph(f"Generado por JARVIS AI - {datetime.now().strftime('%Y-%m-%d %H:%M')}", body_style))
            elements.append(Spacer(1, 12))

            for section in self._expand_content(prompt):
                elements.append(Spacer(1, 8))
                elements.append(Paragraph(section['title'], head_style))
                elements.append(Paragraph(section['body'], body_style))

            doc.build(elements)
            self.log(f"Documento generado: {path}")
            return {"type": "document", "path": path, "prompt": prompt, "format": "PDF"}

        except Exception as e:
            self.log(f"Error generando PDF: {e}")
            # Fallback a TXT
            fname = _safe_name(prompt, 30) + ".txt"
            path = _out("Documentos", fname)
            with open(path, 'w', encoding='utf-8') as f:
                f.write(f"{'='*60}\n  {prompt[:100]}\n  JARVIS AI - {datetime.now().strftime('%Y-%m-%d %H:%M')}\n{'='*60}\n\n")
                for s in self._expand_content(prompt):
                    f.write(f"\n## {s['title']}\n\n{s['body']}\n")
            return {"type": "document", "path": path, "prompt": prompt, "format": "TXT"}

    def _expand_content(self, prompt):
        llm_sections = self._llm_sections(prompt)
        if llm_sections:
            self.log(f"Contenido generado con Llama 3 ({len(llm_sections)} secciones)")
            return llm_sections
        t = prompt.lower()
        if any(k in t for k in ["plan", "proyecto", "roadmap"]):
            return [
                {"title": "Objetivos", "body": "Definir metas claras y alcanzables para el proyecto."},
                {"title": "Fases", "body": "1. Investigacion y analisis. 2. Diseno y planificacion. 3. Desarrollo e implementacion. 4. Pruebas y validacion. 5. Despliegue y mantenimiento."},
                {"title": "Recursos Necesarios", "body": "Equipo humano, herramientas tecnologicas, infraestructura, presupuesto estimado."},
                {"title": "Cronograma", "body": "Estimacion de tiempos por fase: 2-4 semanas cada una."},
                {"title": "Riesgos", "body": "Identificar riesgos potenciales y estrategias de mitigacion."},
            ]
        elif any(k in t for k in ["informe", "reporte", "analisis", "ventas"]):
            return [
                {"title": "Resumen Ejecutivo", "body": "Panorama general de los hallazgos principales."},
                {"title": "Metodologia", "body": "Proceso utilizado para recopilar y analizar la informacion."},
                {"title": "Resultados", "body": "Datos y observaciones clave obtenidos del analisis."},
                {"title": "Conclusiones", "body": "Interpretacion de los resultados y recomendaciones."},
            ]
        elif any(k in t for k in ["contrato", "oficio", "carta"]):
            return [
                {"title": "Encabezado", "body": "Presentacion formal del documento y partes involucradas."},
                {"title": "Terminos y Condiciones", "body": "Detalles de los acuerdos, obligaciones y condiciones establecidas."},
                {"title": "Firma", "body": "Espacios para firmas y sellos de las partes."},
            ]
        else:
            return [
                {"title": "Introduccion", "body": f"Este documento aborda el tema de: {prompt[:80]}."},
                {"title": "Desarrollo", "body": "Analisis detallado del tema solicitado, incluyendo conceptos fundamentales, aplicaciones practicas y consideraciones relevantes."},
                {"title": "Conclusiones", "body": "Resumen de los puntos mas importantes tratados en este documento."},
            ]

    # ═══════════════════════════════════════════════════════════════════════════
    # EXCEL
    # ═══════════════════════════════════════════════════════════════════════════
    def gen_excel(self, prompt):
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

            wb = Workbook()
            ws = wb.active
            ws.title = "Datos"

            hdr_fill = PatternFill(start_color="00c8ff", end_color="00c8ff", fill_type="solid")
            hdr_font = Font(bold=True, color="000000", size=11)
            border = Border(left=Side(style='thin'), right=Side(style='thin'), top=Side(style='thin'), bottom=Side(style='thin'))
            center = Alignment(horizontal='center')

            t = prompt.lower()
            if any(k in t for k in ["ventas", "venta", "commercial"]):
                headers = ["Producto", "Cantidad", "Precio Unit.", "Total", "Fecha"]
                data = [
                    ["Producto A", 150, 25.50, "=B2*C2", "2026-01-15"],
                    ["Producto B", 89, 42.00, "=B3*C3", "2026-01-15"],
                    ["Producto C", 234, 18.75, "=B4*C4", "2026-01-16"],
                    ["Producto D", 67, 55.00, "=B5*C5", "2026-01-16"],
                    ["Producto E", 312, 12.25, "=B6*C6", "2026-01-17"],
                ]
            elif any(k in t for k in ["inventario", "stock", "almacen"]):
                headers = ["Articulo", "Stock Actual", "Stock Minimo", "Estado", "Proveedor"]
                data = [
                    ["Monitor 27\"", 24, 10, "OK", "TechSupply"],
                    ["Teclado Mecanico", 56, 20, "OK", "KeyParts"],
                    ["Mouse Inalambrico", 8, 15, "BAJO", "WirelessCo"],
                    ["Auriculares", 31, 10, "OK", "AudioMax"],
                    ["Webcam HD", 3, 8, "CRITICO", "VisionTech"],
                ]
            elif any(k in t for k in ["nomina", "sueldo", "salario", "empleado"]):
                headers = ["Empleado", "Cargo", "Salario", "Bonos", "Total"]
                data = [
                    ["Juan Perez", "Desarrollador", 3500, 500, "=C2+D2"],
                    ["Maria Garcia", "Disenadora", 3200, 400, "=C3+D3"],
                    ["Carlos Lopez", "Gerente", 5000, 800, "=C4+D4"],
                    ["Ana Martinez", "QA", 2800, 350, "=C5+D5"],
                ]
            else:
                headers = ["Item", "Detalle", "Valor", "Estado", "Notas"]
                data = [
                    ["Elemento 1", "Descripcion del item", 100, "Activo", "Nota 1"],
                    ["Elemento 2", "Descripcion del item", 200, "Pendiente", "Nota 2"],
                    ["Elemento 3", "Descripcion del item", 150, "Completado", "Nota 3"],
                ]

            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=h)
                cell.font = hdr_font
                cell.fill = hdr_fill
                cell.alignment = center
                cell.border = border

            for row_idx, row_data in enumerate(data, 2):
                for col_idx, value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.border = border
                    cell.alignment = center

            for col in ws.columns:
                max_len = max(len(str(cell.value or "")) for cell in col) + 3
                ws.column_dimensions[col[0].column_letter].width = max(max_len, 12)

            fname = _safe_name(prompt, 30) + ".xlsx"
            path = _out("Documentos", fname)
            wb.save(path)
            self.log(f"Excel generado: {path}")
            return {"type": "excel", "path": path, "prompt": prompt}

        except Exception as e:
            self.log(f"Error generando Excel: {e}")
            return {"type": "error", "message": str(e)}

    # ═══════════════════════════════════════════════════════════════════════════
    # POWERPOINT
    # ═══════════════════════════════════════════════════════════════════════════
    def gen_pptx(self, prompt):
        try:
            from pptx import Presentation
            from pptx.util import Inches as In, Pt
            from pptx.dml.color import RGBColor

            prs = Presentation()
            prs.slide_width = In(13.333)
            prs.slide_height = In(7.5)

            # Slide titulo
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x0a, 0x0e, 0x18)
            txBox = slide.shapes.add_textbox(In(1), In(2), In(11), In(2))
            p = txBox.text_frame.paragraphs[0]
            p.text = prompt[:80]
            p.font.size = Pt(36); p.font.color.rgb = RGBColor(0x00, 0xc8, 0xff); p.font.bold = True; p.alignment = 1
            txBox2 = slide.shapes.add_textbox(In(1), In(4.5), In(11), In(1))
            p2 = txBox2.text_frame.paragraphs[0]
            p2.text = f"Generado por JARVIS AI - {datetime.now().strftime('%Y-%m-%d')}"
            p2.font.size = Pt(16); p2.font.color.rgb = RGBColor(0x5a, 0x7a, 0x8a); p2.alignment = 1

            # Slides de contenido - Llama 3 o fallback
            slides_data = []
            llm_slides = self._llm_sections(prompt, 5)
            if llm_slides:
                slides_data = [(s.get("title", f"Slide {i+1}"), s.get("body", "")) for i, s in enumerate(llm_slides)]
                self.log(f"PPTX contenido generado con Llama 3 ({len(slides_data)} slides)")
            else:
                t = prompt.lower()
                if any(k in t for k in ["plan", "proyecto", "roadmap"]):
                    slides_data = [
                        ("Objetivos", "Definir metas y alcance del proyecto\n\n- Meta principal\n- Alcance\n- KPIs de exito"),
                        ("Fases del Proyecto", "1. Investigacion y analisis\n2. Diseno y planificacion\n3. Desarrollo e implementacion\n4. Pruebas y validacion\n5. Despliegue"),
                        ("Recursos", "Equipo humano (5-10 personas)\nHerramientas tecnologicas\nInfraestructura cloud\nPresupuesto estimado"),
                        ("Cronograma", "Fase 1: Semanas 1-2\nFase 2: Semanas 3-5\nFase 3: Semanas 6-10\nFase 4: Semanas 11-12\nFase 5: Semana 13"),
                        ("Riesgos y Mitigacion", "Riesgo 1: Retrasos -> Buffer de tiempo\nRiesgo 2: Presupuesto -> Revision mensual\nRiesgo 3: Tecnologia -> PoC temprano"),
                    ]
                else:
                    slides_data = [
                        ("Introduccion", "Contexto y objetivos de la presentacion"),
                        ("Desarrollo", "Analisis y detalles principales del tema"),
                        ("Resultados", "Hallazgos y datos clave obtenidos"),
                        ("Conclusiones", "Resumen y recomendaciones finales"),
                    ]

            for title, content in slides_data:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x0a, 0x0e, 0x18)
                txBox = slide.shapes.add_textbox(In(1), In(0.5), In(11), In(1.5))
                p = txBox.text_frame.paragraphs[0]
                p.text = title
                p.font.size = Pt(32); p.font.color.rgb = RGBColor(0xff, 0x6a, 0x00); p.font.bold = True
                txBox2 = slide.shapes.add_textbox(In(1), In(2.5), In(11), In(4))
                tf = txBox2.text_frame; tf.word_wrap = True
                p2 = tf.paragraphs[0]
                p2.text = content
                p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(0xc0, 0xd0, 0xe0)

            fname = _safe_name(prompt, 30) + ".pptx"
            path = _out("Documentos", fname)
            prs.save(path)
            self.log(f"PowerPoint generado: {path}")
            return {"type": "pptx", "path": path, "prompt": prompt}

        except Exception as e:
            self.log(f"Error generando PowerPoint: {e}")
            return {"type": "error", "message": str(e)}

    # ═══════════════════════════════════════════════════════════════════════════
    # WORD
    # ═══════════════════════════════════════════════════════════════════════════
    def gen_word(self, prompt):
        try:
            from docx import Document
            from docx.shared import Pt as DocPt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()
            style = doc.styles['Normal']
            style.font.name = 'Calibri'
            style.font.size = DocPt(11)

            title = doc.add_heading(prompt[:80], level=0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in title.runs:
                run.font.color.rgb = RGBColor(0x00, 0xc8, 0xff)

            doc.add_paragraph(f"Generado por JARVIS AI - {datetime.now().strftime('%Y-%m-%d %H:%M')}")

            for section in self._expand_content(prompt):
                h = doc.add_heading(section.get('title', 'Seccion'), level=1)
                for run in h.runs:
                    run.font.color.rgb = RGBColor(0xff, 0x6a, 0x00)
                doc.add_paragraph(section.get('body', ''))

            fname = _safe_name(prompt, 30) + ".docx"
            path = _out("Documentos", fname)
            doc.save(path)
            self.log(f"Word generado: {path}")
            return {"type": "word", "path": path, "prompt": prompt}

        except Exception as e:
            self.log(f"Error generando Word: {e}")
            return {"type": "error", "message": str(e)}

    # ═══════════════════════════════════════════════════════════════════════════
    # CODIGO
    # ═══════════════════════════════════════════════════════════════════════════
    def gen_code(self, prompt):
        t = prompt.lower()
        ts = datetime.now().strftime('%Y-%m-%d %H:%M')

        if any(k in t for k in ["html", "pagina", "web page", "sitio"]):
            ext, code = ".html", f'''<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <title>{prompt[:50]}</title>
    <style>
        body {{ background: #0a0e18; color: #c0d0e0; font-family: Arial, sans-serif; padding: 40px; }}
        h1 {{ color: #00c8ff; }}
        p {{ color: #6a8a9a; line-height: 1.6; }}
    </style>
</head>
<body>
    <h1>{prompt[:50]}</h1>
    <p>Generado por JARVIS AI - {ts}</p>
</body>
</html>'''
        elif any(k in t for k in ["css", "estilo", "style"]):
            ext, code = ".css", f'''/* {prompt} - JARVIS AI {ts} */
* {{ box-sizing: border-box; margin: 0; padding: 0; }}
body {{ background: #0a0e18; color: #c0d0e0; font-family: Arial, sans-serif; }}
.container {{ max-width: 1200px; margin: 0 auto; padding: 20px; }}
'''
        elif any(k in t for k in ["javascript", "js", "react", "node"]):
            ext, code = ".js", f'''// {prompt}
// Generado por JARVIS AI - {ts}

function main() {{
    console.log("=== {prompt[:50]} ===");
}}

main();
'''
        elif any(k in t for k in ["java"]):
            ext, code = ".java", f'''// {prompt}
// Generado por JARVIS AI - {ts}

public class Main {{
    public static void main(String[] args) {{
        System.out.println("=== {prompt[:50]} ===");
    }}
}}
'''
        else:
            ext, code = ".py", f'''#!/usr/bin/env python3
"""
{prompt}
Generado por JARVIS AI - {ts}
"""

def main():
    print("=== {prompt[:50]} ===")

if __name__ == "__main__":
    main()
'''

        fname = _safe_name(prompt, 30) + ext
        path = _out("Codigo", fname)
        with open(path, 'w', encoding='utf-8') as f:
            f.write(code)
        self.log(f"Codigo generado: {path}")
        return {"type": "code", "path": path, "prompt": prompt, "language": ext[1:]}

    # ═══════════════════════════════════════════════════════════════════════════
    # PLAN
    # ═══════════════════════════════════════════════════════════════════════════
    def gen_plan(self, prompt):
        try:
            import matplotlib
            matplotlib.use('Agg')
            import matplotlib.pyplot as plt
            import matplotlib.patches as patches

            fig, ax = plt.subplots(1, 1, figsize=(14, 8), facecolor='#0a0e18')
            ax.set_facecolor('#0a0e18')
            ax.set_xlim(0, 100)
            ax.set_ylim(0, 60)
            ax.axis('off')

            phases = ["Investigacion", "Diseno", "Desarrollo", "Pruebas", "Lanzamiento"]
            colors = ['#00c8ff', '#ff6a00', '#22cc66', '#cc22aa', '#ffcc00']
            durations = [15, 20, 30, 20, 15]

            y = 40
            for i, (phase, color, dur) in enumerate(zip(phases, colors, durations)):
                x = sum(durations[:i]) + 5
                rect = patches.FancyBboxPatch((x, y-4), dur-2, 8, boxstyle="round,pad=0.4",
                                               facecolor='#0d1520', edgecolor=color, linewidth=2.5)
                ax.add_patch(rect)
                ax.text(x + (dur-2)/2, y+1, phase, ha='center', va='center', color=color, fontsize=11, fontweight='bold')
                ax.text(x + (dur-2)/2, y-6, f"Sem {sum(durations[:i])+1}-{sum(durations[:i+1])}", ha='center', va='center', color='#4a6878', fontsize=8)

            ax.text(50, 55, prompt[:60], ha='center', va='center', color='#00c8ff', fontsize=15, fontweight='bold')

            fname = _safe_name(prompt, 30) + "_plan.png"
            path = _out("Planes", fname)
            fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#0a0e18')
            plt.close(fig)
            self.log(f"Plan generado: {path}")
            return {"type": "plan", "path": path, "prompt": prompt}

        except Exception as e:
            self.log(f"Error generando plan: {e}")
            return {"type": "error", "message": str(e)}

    # ═══════════════════════════════════════════════════════════════════════════
    # MUSICA
    # ═══════════════════════════════════════════════════════════════════════════
    def gen_music(self, prompt):
        try:
            import struct, wave as wave_mod

            sample_rate = 44100
            duration = 8
            t = prompt.lower()
            if any(k in t for k in ["relaj", "chill", "calm", "meditacion"]):
                base_freq = 220
                dur = 12
            elif any(k in t for k in ["ener", "rock", "fast", "dance"]):
                base_freq = 440
                dur = 6
            else:
                base_freq = 330
                dur = 8

            fname = _safe_name(prompt, 30) + ".wav"
            path = _out("Audio", fname)

            with wave_mod.open(path, 'w') as wav:
                wav.setnchannels(1)
                wav.setsampwidth(2)
                wav.setframerate(sample_rate)
                frames = []
                for i in range(sample_rate * dur):
                    t_norm = i / sample_rate
                    # Melodia con armonicos
                    sample = 0
                    sample += 12000 * math.sin(2 * math.pi * base_freq * t_norm)
                    sample += 6000 * math.sin(2 * math.pi * base_freq * 1.5 * t_norm)
                    sample += 3000 * math.sin(2 * math.pi * base_freq * 2 * t_norm)
                    # Envolvente
                    env = max(0, 1 - t_norm / dur) * min(1, t_norm * 4)
                    sample *= env
                    frames.append(struct.pack('<h', max(-32768, min(32767, int(sample)))))
                wav.writeframes(b''.join(frames))

            self.log(f"Audio generado: {path}")
            return {"type": "audio", "path": path, "prompt": prompt, "format": "WAV"}

        except Exception as e:
            self.log(f"Error generando audio: {e}")
            return {"type": "error", "message": str(e)}
